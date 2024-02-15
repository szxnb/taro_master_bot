from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime

def generate_report_ppt(part1, part2, part3, part4):
    print("开始生成PPT")
    # 创建一个新的PPT对象
    prs = Presentation()

    slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
    title_1 = slide_1.shapes.title
    title_1.text = "塔罗大师报告"

    slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
    title_2 = slide_2.shapes.title
    title_2.text = "总述"
    content_2 = slide_2.placeholders[1].text = part1

    slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
    title_3 = slide_3.shapes.title
    title_3.text = "感情方面测算"
    content_3 = slide_3.placeholders[1].text = part2

    slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
    title_4 = slide_4.shapes.title
    title_4.text = "职业方面测算"
    content_4 = slide_4.placeholders[1].text = part3

    slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
    title_5 = slide_5.shapes.title
    title_5.text = "财务方面测算"
    content_5 = slide_5.placeholders[1].text = part4

    # 设置文本框的字体大小为 10
    font_size = Pt(20)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = font_size

    # 保存PPT文件
    # 格式化当前时间
    formatted_time = datetime.now().strftime("%Y%m%d%H%M%S")
    report_name = "reports/TaroReport_" + "3500466989@qq.com_" + formatted_time + ".pptx"
    print(report_name)
    prs.save("./" + report_name)
    print("生成PPT完成!")