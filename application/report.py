from datetime import datetime
from pptx import Presentation
from card_detail import card_detail
import os
import subprocess

def replace_content(shape, target_text, replacement_text):
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if target_text in run.text:
                    run.text = run.text.replace(target_text, replacement_text)

def find_image_shape(slide):
    num = 0
    # 遍历幻灯片中的所有形状
    for shape in slide.shapes:
        # 检查形状是否为图片
        if shape.shape_type == 13:  # 13
            num = num + 1
            if num == 2:
                return shape
    return None

def replace_image(slide, new_image_path):
    # 查找包含图片的形状
    old_image_shape = find_image_shape(slide)

    if old_image_shape is not None:
        # 记录旧图片的位置信息
        left = old_image_shape.left
        top = old_image_shape.top
        width = old_image_shape.width
        height = old_image_shape.height

        # 删除旧图片
        slide.shapes._spTree.remove(old_image_shape._element)

        # 插入新图片
        slide.shapes.add_picture(new_image_path, left, top, width, height)

def generate_ppt(part1, part2, part3, part4):
    presentation = Presentation("../templates/template.pptx")
    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            replace_content(shape, "{{#1overview}}", part1)
            replace_content(shape, "{{#2love}}", part2)
            replace_content(shape, "{{#3career}}", part3)
            replace_content(shape, "{{#4finance}}", part4)
    # 格式化当前时间
    formatted_time = datetime.now().strftime("%Y%m%d%H%M%S")
    report_name = "reports/TaroReport_" + "3500466989@qq.com_" + formatted_time + ".pptx"
    print(report_name)
    presentation.save("../" + report_name)
    print("生成PPT完成!")


def generate_pdf(qusetion=None, card1=None, card2=None, card3=None, answers=None, recipient_email=None):
    # 获取当前脚本所在的目录
    current_script_directory = os.path.dirname(os.path.abspath(__file__))
    # 获取项目根目录
    project_root = os.path.abspath(os.path.join(current_script_directory, ".."))
    presentation = Presentation(os.path.join(project_root, "templates/new_template.pptx"))

    user_question = qusetion
    master_words = 'Before we begin the interpretation, I would like to explain some important points to you. Tarot cards are not tools for predicting the future, but rather tools for guiding us to delve deeper into various situations in life. Therefore, when receiving the reading, please maintain an open mindset and trust your intuition. Our goal is to provide you with inspiration and guidance through the guidance of the tarot cards, helping you better deal with challenges and decisions in life. Alright, you can start reading now!'
    # Name
    card_first_name = 'Card Ⅰ ' + card_detail[card1][0]
    card_second_name = 'Card Ⅱ ' + card_detail[card2][0]
    card_third_name = 'Card Ⅲ ' + card_detail[card3][0]
    # Image
    card_first_image = os.path.join(project_root, "assets\\", card_detail[card1][1])
    card_second_image = os.path.join(project_root, "assets\\", card_detail[card2][1])
    card_third_image = os.path.join(project_root, "assets\\", card_detail[card3][1])

    print( card_first_image,card_second_image, card_third_image)
    # Description
    card_first_description = card_detail[card1][2]
    card_second_description = card_detail[card2][2]
    card_third_description = card_detail[card3][2]

    summarize = answers['result_overview']
    love_suggestion = answers['result_love']
    career_suggestion = answers['result_career']
    wealth_suggestion = answers['result_finances']


    # 获取当前时间
    current_time = datetime.now()
    # 将时间格式化为 "Month day, year" 格式
    formatted_time = current_time.strftime("%B %d, %Y")
    for slide_number, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            replace_content(shape, "{{#date}}", formatted_time)
            replace_content(shape, "{{#user_question}}", user_question)
            replace_content(shape, "{{#master_words}}", master_words)
            replace_content(shape, "{{#card_first_name}}", card_first_name)
            replace_content(shape, "{{#card_first_description}}", card_first_description)
            replace_content(shape, "{{#card_second_name}}", card_second_name)
            replace_content(shape, "{{#card_second_description}}", card_second_description)
            replace_content(shape, "{{#card_third_name}}", card_third_name)
            replace_content(shape, "{{#card_third_description}}", card_third_description)
            replace_content(shape, "{{#summarize}}", summarize)
            replace_content(shape, "{{#love_suggestion}}", love_suggestion)
            replace_content(shape, "{{#career_suggestion}}", career_suggestion)
            replace_content(shape, "{{#wealth_suggestion}}", wealth_suggestion)
    slide = presentation.slides[2]
    replace_image(slide, card_first_image)
    slide = presentation.slides[3]
    replace_image(slide, card_second_image)
    slide = presentation.slides[4]
    replace_image(slide, card_third_image)
    # 格式化当前时间
    formatted_time = datetime.now().strftime("%Y%m%d%H%M%S")
    report_name = "reports\\TaroReport_" + recipient_email +"_" + formatted_time + ".pptx"
    presentation.save("../" + report_name)
    print("生成PPT完成!")
    pptx_path = os.path.join(project_root, report_name)
    # 定义命令
    command = [
        'soffice',
        '--headless',
        '--convert-to',
        'pdf:impress_pdf_Export',
        '--outdir',
        os.path.join(project_root, "reports"),
        pptx_path
    ]
    # 执行命令
    process = subprocess.run(command, capture_output=True, text=True)
    # 打印输出结果和错误信息
    print("Command Output:", process.stdout)
    print("Command Error:", process.stderr)
    # 打印返回码
    print("Return Code:", process.returncode)
    return pptx_path.replace(".pptx",".pdf")


if __name__ == "__main__":
    generate_pdf("", "Fool", "Magician", "Priestess")
